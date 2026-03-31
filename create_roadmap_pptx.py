"""
UAE PASS Digital Documents 2026 Roadmap - PowerPoint Generator
Creates a professional presentation with color-coded categories and quarterly timeline
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color scheme (using RGBColor instead of RgbColor)
COLORS = {
    'product': RGBColor(0x00, 0x7A, 0xCC),      # Blue
    'design': RGBColor(0x8B, 0x5C, 0xF6),       # Purple
    'sp': RGBColor(0x14, 0xB8, 0xA6),           # Teal
    'ux': RGBColor(0xF9, 0x73, 0x16),           # Orange
    'title_bg': RGBColor(0x1E, 0x3A, 0x5F),     # Dark Blue
    'white': RGBColor(0xFF, 0xFF, 0xFF),
    'light_gray': RGBColor(0xF5, 0xF5, 0xF5),
    'dark_gray': RGBColor(0x33, 0x33, 0x33),
    'medium_gray': RGBColor(0x66, 0x66, 0x66),
}

def add_title_slide(prs):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background shape (dark blue)
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['title_bg']
    bg.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "UAE PASS Digital Documents"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.8), Inches(9), Inches(1)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Product Roadmap"
    p.font.size = Pt(32)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Decorative line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(3.5), Inches(4.6), Inches(3), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['product']
    line.line.fill.background()

    # Footer text
    footer_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Empowering UAE PASS Users and Service Providers"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xA0, 0xA0, 0xA0)
    p.alignment = PP_ALIGN.CENTER

def add_overview_slide(prs):
    """Create the roadmap overview slide with quarterly timeline"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Slide title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Roadmap Overview"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    # Quarter headers
    quarters = ['Q1 2026', 'Q2 2026', 'Q3 2026', 'Q4 2026']
    q_width = Inches(2.2)
    q_start = Inches(0.4)
    q_top = Inches(0.9)

    for i, q in enumerate(quarters):
        q_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            q_start + (i * q_width) + Inches(i * 0.1),
            q_top,
            q_width,
            Inches(0.4)
        )
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = COLORS['title_bg']
        q_box.line.fill.background()

        tf = q_box.text_frame
        tf.paragraphs[0].text = q
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False

    # Category rows with features
    categories = [
        ('PRODUCT', COLORS['product'], [
            ('Download All Documents', 'Q1'),
            ('Auto-Add Documents', 'Q1-Q2'),
            ('Status Reporting', 'Q1'),
            ('Analytics Tool', 'Q2'),
            ('Dual Citizenship GA', 'Q1'),
            ('Consent Sharing', 'Q3-Q4'),
            ('Physical Doc Sharing', 'Q4'),
        ]),
        ('DESIGN', COLORS['design'], [
            ('Design Audit', 'Q1'),
            ('Design System Update', 'Q2'),
            ('Accessibility', 'Q2-Q3'),
        ]),
        ('SERVICE PROVIDER', COLORS['sp'], [
            ('Form Filler', 'Q1-Q2'),
        ]),
        ('UX', COLORS['ux'], [
            ('Home Page Revamp', 'Q2'),
            ('Documents List View', 'Q1-Q2'),
            ('UX Enhancements', 'Q2-Q3'),
            ('Document Request Flow', 'Q3'),
            ('Document Sharing Flow', 'Q3-Q4'),
        ]),
    ]

    row_top = Inches(1.5)
    row_height = Inches(1.2)
    label_width = Inches(1.2)

    quarter_map = {
        'Q1': (0, 1),
        'Q2': (1, 1),
        'Q3': (2, 1),
        'Q4': (3, 1),
        'Q1-Q2': (0, 2),
        'Q2-Q3': (1, 2),
        'Q3-Q4': (2, 2),
        'Q1-Q3': (0, 3),
        'Q2-Q4': (1, 3),
    }

    for cat_idx, (cat_name, cat_color, features) in enumerate(categories):
        current_top = row_top + (cat_idx * row_height)

        # Category label
        cat_label = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.15),
            current_top,
            label_width,
            row_height - Inches(0.15)
        )
        cat_label.fill.solid()
        cat_label.fill.fore_color.rgb = cat_color
        cat_label.line.fill.background()

        tf = cat_label.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = cat_name
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = COLORS['white']
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Feature boxes in timeline
        feat_height = Inches(0.28)
        feat_gap = Inches(0.05)

        row_features = []
        for feat_name, timing in features:
            start_q, span = quarter_map.get(timing, (0, 1))
            row_features.append((feat_name, start_q, span))

        for feat_idx, (feat_name, start_q, span) in enumerate(row_features):
            feat_left = q_start + (start_q * q_width) + (start_q * Inches(0.1)) + Inches(0.05)
            feat_width = (span * q_width) + ((span - 1) * Inches(0.1)) - Inches(0.1)

            # Vertical position within category row
            feat_top = current_top + Inches(0.1) + (feat_idx * (feat_height + feat_gap))

            if feat_top + feat_height > current_top + row_height - Inches(0.1):
                continue

            feat_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                feat_left,
                feat_top,
                feat_width,
                feat_height
            )
            feat_box.fill.solid()
            feat_box.fill.fore_color.rgb = cat_color
            feat_box.line.fill.background()

            tf = feat_box.text_frame
            tf.word_wrap = False
            tf.paragraphs[0].text = feat_name
            tf.paragraphs[0].font.size = Pt(8)
            tf.paragraphs[0].font.color.rgb = COLORS['white']
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Legend at bottom
    legend_top = Inches(6.3)
    legend_items = [
        ('Product', COLORS['product']),
        ('Design', COLORS['design']),
        ('Service Provider', COLORS['sp']),
        ('UX', COLORS['ux']),
    ]

    legend_start = Inches(2)
    for i, (label, color) in enumerate(legend_items):
        # Color box
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            legend_start + (i * Inches(1.8)),
            legend_top,
            Inches(0.25),
            Inches(0.25)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        # Label
        lbl = slide.shapes.add_textbox(
            legend_start + (i * Inches(1.8)) + Inches(0.3),
            legend_top,
            Inches(1.4),
            Inches(0.3)
        )
        tf = lbl.text_frame
        tf.paragraphs[0].text = label
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = COLORS['dark_gray']

def add_category_header(slide, title, color, subtitle=None):
    """Add a category header to a slide"""
    # Header bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, Inches(10), Inches(1.1)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = color
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.25), Inches(9), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.7), Inches(9), Inches(0.4)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)

def add_feature_box(slide, left, top, width, height, title, description, quarter, audience, color):
    """Add a feature box to a slide"""
    # Main box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLORS['light_gray']
    box.line.color.rgb = color
    box.line.width = Pt(2)

    # Quarter badge
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left + width - Inches(0.8),
        top + Inches(0.1),
        Inches(0.7),
        Inches(0.25)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()

    tf = badge.text_frame
    tf.paragraphs[0].text = quarter
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['white']
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + Inches(0.1),
        width - Inches(1),
        Inches(0.35)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    # Description
    desc_box = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + Inches(0.45),
        width - Inches(0.3),
        height - Inches(0.85)
    )
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = description
    p.font.size = Pt(10)
    p.font.color.rgb = COLORS['medium_gray']

    # Audience label
    aud_box = slide.shapes.add_textbox(
        left + Inches(0.15),
        top + height - Inches(0.35),
        width - Inches(0.3),
        Inches(0.3)
    )
    tf = aud_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"For: {audience}"
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = color

def add_product_slides(prs):
    """Create Product category slides"""
    slide_layout = prs.slide_layouts[6]

    # Slide 1: Core Document Features
    slide = prs.slides.add_slide(slide_layout)
    add_category_header(slide, "Product Features", COLORS['product'], "Core Document Capabilities")

    features = [
        ("Download All Documents",
         "Enable users to download all their issued documents from DV for offline access and personal records.",
         "Q1", "UAE PASS Users"),
        ("Auto-Add Documents",
         "One-time consent allowing DV to periodically check and auto-add new or updated documents from issuers. Streamlines document collection.",
         "Q1-Q2", "UAE PASS Users"),
        ("Dual Citizenship GA",
         "Complete Primary/Secondary EID handling for dual citizens, with full EN/AR support and edge case handling.",
         "Q1", "UAE PASS Users"),
    ]

    col_width = Inches(3)
    gap = Inches(0.15)
    start_left = Inches(0.35)
    top = Inches(1.4)
    height = Inches(1.6)

    for i, (title, desc, quarter, audience) in enumerate(features):
        left = start_left + (i * (col_width + gap))
        add_feature_box(slide, left, top, col_width, height, title, desc, quarter, audience, COLORS['product'])

    # Slide 2: Analytics & Data Features
    slide = prs.slides.add_slide(slide_layout)
    add_category_header(slide, "Product Features", COLORS['product'], "Analytics & Insights")

    features = [
        ("Status-Based Reporting",
         "Deploy 23-status-code tracking system for accurate measurement of all sharing requests. Foundation for optimization and insights.",
         "Q1", "Service Providers"),
        ("User Behavior Analytics",
         "Select and integrate analytics tool (UXCam/Firebase) to enable data-driven optimization, session replay, and Real User Monitoring.",
         "Q2", "UAE PASS Users"),
    ]

    col_width = Inches(4.5)
    for i, (title, desc, quarter, audience) in enumerate(features):
        left = start_left + (i * (col_width + gap))
        add_feature_box(slide, left, top, col_width, Inches(1.5), title, desc, quarter, audience, COLORS['product'])

    # Slide 3: Advanced Sharing Features
    slide = prs.slides.add_slide(slide_layout)
    add_category_header(slide, "Product Features", COLORS['product'], "Advanced Sharing Capabilities")

    features = [
        ("Consent Sharing (Third Party)",
         "Enable sharing of third-party data through a robust consent mechanism. Expands data sharing possibilities while maintaining user control.",
         "Q3-Q4", "Service Providers"),
        ("Physical Document Sharing",
         "Enable sharing of physical document representations or verified copies for scenarios requiring tangible documentation.",
         "Q4", "Service Providers"),
    ]

    col_width = Inches(4.5)
    for i, (title, desc, quarter, audience) in enumerate(features):
        left = start_left + (i * (col_width + gap))
        add_feature_box(slide, left, top, col_width, Inches(1.5), title, desc, quarter, audience, COLORS['product'])

def add_design_slide(prs):
    """Create Design category slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_category_header(slide, "Design Initiatives", COLORS['design'], "Visual & Accessibility Improvements")

    features = [
        ("Design Audit",
         "Comprehensive audit of current design patterns to identify improvement areas and establish baseline metrics.",
         "Q1", "UAE PASS Users"),
        ("Design System Update",
         "Update and maintain the DV design system for consistency across all platforms and touchpoints.",
         "Q2", "UAE PASS Users"),
        ("Accessibility Enhancement",
         "Improve accessibility compliance (WCAG) across the application to ensure inclusive access for all users.",
         "Q2-Q3", "UAE PASS Users"),
    ]

    col_width = Inches(3)
    gap = Inches(0.15)
    start_left = Inches(0.35)
    top = Inches(1.4)
    height = Inches(1.6)

    for i, (title, desc, quarter, audience) in enumerate(features):
        left = start_left + (i * (col_width + gap))
        add_feature_box(slide, left, top, col_width, height, title, desc, quarter, audience, COLORS['design'])

def add_sp_slide(prs):
    """Create Service Provider category slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_category_header(slide, "Service Provider Features", COLORS['sp'], "Integration Capabilities")

    # Single feature, centered
    title = "Form Filler"
    desc = "Auto-fill forms using stored document data to streamline user workflows. Reduces manual data entry and errors while improving completion rates for service applications."
    quarter = "Q1-Q2"
    audience = "Service Providers, UAE PASS Users"

    add_feature_box(slide, Inches(2.5), Inches(1.8), Inches(5), Inches(1.8),
                    title, desc, quarter, audience, COLORS['sp'])

    # Additional context box
    info_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(4), Inches(7), Inches(1.5)
    )
    tf = info_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Key Benefits"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_gray']

    benefits = [
        "Reduced form completion time for users",
        "Higher conversion rates for service providers",
        "Decreased data entry errors",
        "Seamless integration with existing SP systems"
    ]

    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = f"  {benefit}"
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['medium_gray']
        p.space_before = Pt(4)

def add_ux_slides(prs):
    """Create UX category slides"""
    slide_layout = prs.slide_layouts[6]

    # Slide 1: Core UX Improvements
    slide = prs.slides.add_slide(slide_layout)
    add_category_header(slide, "UX Enhancements", COLORS['ux'], "User Experience Improvements")

    features = [
        ("Home Page Revamp",
         "Redesign the DV home page for improved user experience, engagement, and easier access to key features.",
         "Q2", "UAE PASS Users"),
        ("Documents List View",
         "UI/UX enhancements for documents list view including improved navigation, filtering, and visual hierarchy.",
         "Q1-Q2", "UAE PASS Users"),
        ("UX Enhancements Bundle",
         "Additional UI/UX improvements based on user feedback and analytics insights to optimize the overall experience.",
         "Q2-Q3", "UAE PASS Users"),
    ]

    col_width = Inches(3)
    gap = Inches(0.15)
    start_left = Inches(0.35)
    top = Inches(1.4)
    height = Inches(1.6)

    for i, (title, desc, quarter, audience) in enumerate(features):
        left = start_left + (i * (col_width + gap))
        add_feature_box(slide, left, top, col_width, height, title, desc, quarter, audience, COLORS['ux'])

    # Slide 2: Flow Improvements
    slide = prs.slides.add_slide(slide_layout)
    add_category_header(slide, "UX Enhancements", COLORS['ux'], "Critical Flow Redesigns")

    features = [
        ("Document Request Flow",
         "Redesign the document request flow for improved clarity and reduced friction. Focus on guiding users through the process efficiently.",
         "Q3", "UAE PASS Users"),
        ("Document Sharing Flow",
         "Redesign document sharing flow to increase success rates and reduce drop-offs. Address consent screen and flow optimization.",
         "Q3-Q4", "UAE PASS Users, Service Providers"),
    ]

    col_width = Inches(4.5)
    for i, (title, desc, quarter, audience) in enumerate(features):
        left = start_left + (i * (col_width + gap))
        add_feature_box(slide, left, top, col_width, Inches(1.6), title, desc, quarter, audience, COLORS['ux'])

def add_summary_slide(prs):
    """Create summary/closing slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['title_bg']
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "2026 Roadmap Summary"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Stats row
    stats = [
        ("16", "Features"),
        ("4", "Categories"),
        ("4", "Quarters"),
    ]

    stat_width = Inches(2.5)
    stat_start = Inches(1.25)
    stat_top = Inches(2.8)

    for i, (number, label) in enumerate(stats):
        # Number
        num_box = slide.shapes.add_textbox(
            stat_start + (i * stat_width),
            stat_top,
            stat_width,
            Inches(0.8)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = COLORS['product']
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(
            stat_start + (i * stat_width),
            stat_top + Inches(0.7),
            stat_width,
            Inches(0.4)
        )
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    # Key focus areas
    focus_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.3), Inches(8), Inches(2)
    )
    tf = focus_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Key Focus Areas for 2026"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    focuses = [
        "Enhanced document management and accessibility",
        "Improved sharing flows and conversion rates",
        "Comprehensive design system and accessibility compliance",
        "Data-driven optimization with analytics"
    ]

    for focus in focuses:
        p = tf.add_paragraph()
        p.text = focus
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xC0, 0xC0, 0xC0)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

def main():
    """Generate the PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create slides
    add_title_slide(prs)
    add_overview_slide(prs)
    add_product_slides(prs)
    add_design_slide(prs)
    add_sp_slide(prs)
    add_ux_slides(prs)
    add_summary_slide(prs)

    # Save
    output_path = r"D:\claude\UAE_PASS_DV_2026_Roadmap.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
