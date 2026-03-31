"""
Create 'Ask Your Documents' AI Feature Presentation for UAE PASS Digital Vault.
Version 2 — Refined design with 11 slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Brand Colors --
DARK_BLUE = RGBColor(0x00, 0x2D, 0x62)      # UAE PASS primary
DARK_BLUE_2 = RGBColor(0x00, 0x3A, 0x7A)    # Slightly lighter dark blue
MEDIUM_BLUE = RGBColor(0x00, 0x5C, 0x99)
LIGHT_BLUE = RGBColor(0xE8, 0xF1, 0xFA)
GOLD = RGBColor(0xC8, 0xA2, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MEDIUM_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x00, 0x7A, 0x33)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
TEAL = RGBColor(0x00, 0x96, 0x88)
PURPLE = RGBColor(0x8E, 0x44, 0xAD)

TOTAL_SLIDES = 11

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text_box(slide, left, top, width, height, lines, font_size=18,
                           color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                           font_name="Calibri", line_spacing=1.2):
    """Add a text box with multiple paragraphs, each with potentially different formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, l_color, l_bold, l_size = line_info, color, bold, font_size
        else:
            text = line_info[0]
            l_color = line_info[1] if len(line_info) > 1 else color
            l_bold = line_info[2] if len(line_info) > 2 else bold
            l_size = line_info[3] if len(line_info) > 3 else font_size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(l_size)
        p.font.color.rgb = l_color
        p.font.bold = l_bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(2)
        if line_spacing != 1.0:
            p.line_spacing = Pt(l_size * line_spacing)
    return txBox


def add_gold_accent_bar(slide):
    add_rect(slide, Inches(0), H - Inches(0.1), W, Inches(0.1), GOLD)


def add_slide_number(slide, num, total=TOTAL_SLIDES):
    add_text_box(slide, W - Inches(1.2), H - Inches(0.5), Inches(1), Inches(0.4),
                 f"{num} / {total}", font_size=10, color=MEDIUM_GRAY,
                 alignment=PP_ALIGN.RIGHT)


def add_header_bar(slide, title_text, subtitle_text=None):
    """Standard header bar for body slides."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.3), DARK_BLUE)
    add_text_box(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.7),
                 title_text, font_size=32, color=WHITE, bold=True)
    if subtitle_text:
        add_text_box(slide, Inches(0.8), Inches(0.75), Inches(10), Inches(0.4),
                     subtitle_text, font_size=16, color=GOLD)


# ===================================================================
# SLIDE 1: Title Slide — Full Redesign
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, DARK_BLUE)

# Layered gradient-like rectangles (dark to slightly lighter)
add_rect(slide, Inches(0), Inches(0), W, Inches(2.5), RGBColor(0x00, 0x22, 0x4E))
add_rect(slide, Inches(0), Inches(2.5), W, Inches(2.5), DARK_BLUE)
add_rect(slide, Inches(0), Inches(5.0), W, Inches(2.5), DARK_BLUE_2)

# Large decorative translucent circle (arc effect) - right side
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1.0), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x00, 0x3A, 0x7A)
circle.line.fill.background()

# Second smaller decorative circle overlapping
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(0.5), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x00, 0x45, 0x8B)
circle2.line.fill.background()

# Thin gold line at top
add_rect(slide, Inches(0), Inches(0), W, Inches(0.04), GOLD)

# Top label: "UAE PASS Digital Vault"
add_text_box(slide, Inches(1.2), Inches(1.0), Inches(6), Inches(0.5),
             "UAE PASS Digital Vault", font_size=16, color=GOLD, bold=False,
             font_name="Calibri Light")

# Main title
add_text_box(slide, Inches(1.2), Inches(1.8), Inches(8), Inches(1.2),
             '"Ask Your Documents"', font_size=52, color=WHITE, bold=True)

# Subtitle — refined tagline
add_text_box(slide, Inches(1.2), Inches(3.1), Inches(8), Inches(0.8),
             "Your Documents. Your Questions. Instant Answers.",
             font_size=26, color=RGBColor(0xCC, 0xDD, 0xEE), bold=False,
             font_name="Calibri Light")

# Thin gold separator line
add_rect(slide, Inches(1.2), Inches(4.1), Inches(3.0), Inches(0.025), GOLD)

# Metadata line
add_text_box(slide, Inches(1.2), Inches(4.4), Inches(8), Inches(0.5),
             "AI-Powered Document Intelligence  |  Product Management, Digital Vault",
             font_size=14, color=RGBColor(0x88, 0xAA, 0xCC), bold=False,
             font_name="Calibri Light")

# Bottom thin gold bar
add_rect(slide, Inches(0), H - Inches(0.04), W, Inches(0.04), GOLD)

add_slide_number(slide, 1)


# ===================================================================
# SLIDE 2: The Problem
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "The Problem", "Why users need intelligent document access")
add_gold_accent_bar(slide)

problems = [
    ("Information Overload", "Users accumulate 10+ documents\nin their Digital Vault over time.\nFinding specific data requires opening\neach document individually."),
    ("Time-Consuming Search", "Manual browsing through PDFs to find\na single data point (expiry date, ID number,\ncontract clause) wastes valuable time."),
    ("Keyword-Only Search", "Current search is limited to document titles\nand metadata. It cannot search inside\ndocument content or understand context."),
    ("Cross-Document Blind Spots", "No way to ask questions spanning multiple\ndocuments, like 'What expires next month?'\nor 'Summarize my identity details.'"),
]

card_w = Inches(5.4)
card_h = Inches(2.3)
start_x = Inches(0.8)
start_y = Inches(1.7)
gap_x = Inches(0.9)
gap_y = Inches(0.5)

for i, (title, desc) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    card = add_rect(slide, x, y, card_w, card_h, LIGHT_BLUE, MEDIUM_BLUE)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.25), card_w - Inches(0.6), Inches(0.5),
                 title, font_size=20, color=DARK_BLUE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), card_w - Inches(0.6), Inches(1.3),
                 desc, font_size=14, color=MEDIUM_GRAY, line_spacing=1.4)

add_slide_number(slide, 2)


# ===================================================================
# SLIDE 3: How You Interact (NEW)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "How You Interact", "Multiple ways to query your documents")
add_gold_accent_bar(slide)

interaction_methods = [
    ("Text Chat", MEDIUM_BLUE,
     "Type your question naturally in English or Arabic. The AI understands context and returns cited answers from your documents.",
     "\u2328",  # keyboard icon placeholder
     '"When does my visa expire?"  \u00AB\u0645\u062A\u0649 \u062A\u0646\u062A\u0647\u064A \u0635\u0644\u0627\u062D\u064A\u0629 \u062A\u0623\u0634\u064A\u0631\u062A\u064A\u061F\u00BB'),
    ("Voice Input", GREEN,
     "Tap the microphone and ask your question out loud. Speech-to-text handles both English and Arabic seamlessly.",
     "\U0001F3A4",  # microphone
     "Supports natural speech in EN/AR\nwith real-time transcription"),
    ("Quick Actions", GOLD,
     'Pre-built shortcuts for common queries. One tap to get instant answers without typing.',
     "\u26A1",  # lightning bolt
     '"What expires soon?"  \u2022  "My identity summary"\n"Show all contracts"  \u2022  "Renewal reminders"'),
]

card_w = Inches(3.6)
card_h = Inches(5.0)
start_x = Inches(0.8)
card_y = Inches(1.7)
gap = Inches(0.45)

for i, (title, accent, desc, icon, example) in enumerate(interaction_methods):
    x = start_x + i * (card_w + gap)

    # Card
    add_rounded_rect(slide, x, card_y, card_w, card_h, LIGHT_GRAY)
    # Top accent bar
    add_rect(slide, x + Inches(0.02), card_y, card_w - Inches(0.04), Inches(0.07), accent)

    # Icon circle
    icon_size = Inches(0.9)
    icon_x = x + (card_w - icon_size) / 2
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, icon_x, card_y + Inches(0.4), icon_size, icon_size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_text_box(slide, x + Inches(0.2), card_y + Inches(1.5), card_w - Inches(0.4), Inches(0.5),
                 title, font_size=22, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Inches(0.25), card_y + Inches(2.1), card_w - Inches(0.5), Inches(1.3),
                 desc, font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

    # Example area
    example_y = card_y + Inches(3.4)
    add_rounded_rect(slide, x + Inches(0.15), example_y, card_w - Inches(0.3), Inches(1.3), WHITE, RGBColor(0xDD, 0xDD, 0xDD))
    add_text_box(slide, x + Inches(0.25), example_y + Inches(0.1), card_w - Inches(0.5), Inches(1.1),
                 example, font_size=11, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.5)

add_slide_number(slide, 3)


# ===================================================================
# SLIDE 4: The Solution (was Slide 3)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, 'The Solution \u2014 "Ask Your Documents"',
               "Natural language queries across your entire document collection")
add_gold_accent_bar(slide)

features = [
    ("Natural Language\nQueries", "Ask questions in plain English or Arabic.\nNo keywords needed \u2014 just speak naturally.\n\n\"When does my visa expire?\"\n\u00AB\u0645\u062A\u0649 \u062A\u0646\u062A\u0647\u064A \u0635\u0644\u0627\u062D\u064A\u0629 \u062A\u0623\u0634\u064A\u0631\u062A\u064A\u061F\u00BB", MEDIUM_BLUE),
    ("Intelligent\nRetrieval", "AI searches across ALL your documents\nto find the most relevant information.\nWorks with issued and uploaded docs.", GREEN),
    ("Cited\nAnswers", "Every answer includes a citation to the\nsource document, so you can verify.\nNo unsupported claims \u2014 fully grounded.", GOLD),
    ("Privacy-First\nDesign", "Documents NEVER leave your device.\nOnly your question + tiny excerpts\nare sent to a UAE-region cloud for answers.", DARK_BLUE),
]

for i, (title, desc, accent) in enumerate(features):
    x = Inches(0.8) + i * Inches(3.05)
    y = Inches(1.8)
    w = Inches(2.8)
    h = Inches(4.8)

    # Card background
    card = add_rounded_rect(slide, x, y, w, h, LIGHT_GRAY)
    # Top accent
    add_rect(slide, x + Inches(0.02), y, w - Inches(0.04), Inches(0.06), accent)

    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), w - Inches(0.4), Inches(0.7),
                 title, font_size=18, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.2)

    add_text_box(slide, x + Inches(0.2), y + Inches(1.2), w - Inches(0.4), Inches(3.3),
                 desc, font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.5)

add_slide_number(slide, 4)


# ===================================================================
# SLIDE 5: Key Use Cases — Conversational Mockups (was Slide 4)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Key Use Cases", "Conversational AI in action")
add_gold_accent_bar(slide)

# --- Use Case 1: Renewal Planning Assistant ---
uc1_x = Inches(0.6)
uc1_w = Inches(5.8)
uc1_y = Inches(1.55)

add_text_box(slide, uc1_x, uc1_y, uc1_w, Inches(0.45),
             "Renewal Planning Assistant", font_size=18, color=DARK_BLUE, bold=True)

# Chat conversation mockup
chat_y = uc1_y + Inches(0.55)
bubble_gap = Inches(0.15)

# User bubble 1
ub1 = add_rounded_rect(slide, uc1_x, chat_y, Inches(4.2), Inches(0.5), LIGHT_BLUE, MEDIUM_BLUE)
add_text_box(slide, uc1_x + Inches(0.15), chat_y + Inches(0.08), Inches(3.9), Inches(0.35),
             "\U0001F464  \"What documents expire in the next 6 months?\"",
             font_size=12, color=DARK_BLUE, bold=False)

# AI response 1
ai1_y = chat_y + Inches(0.6) + bubble_gap
ai1_h = Inches(1.5)
ai1 = add_rounded_rect(slide, uc1_x + Inches(0.5), ai1_y, Inches(5.1), ai1_h, RGBColor(0xF0, 0xF8, 0xF0), GREEN)
add_multiline_text_box(slide, uc1_x + Inches(0.65), ai1_y + Inches(0.08), Inches(4.8), ai1_h - Inches(0.1),
    [
        ("\U0001F916  Based on your documents, 3 items expire soon:", GREEN, True, 12),
        ("", MEDIUM_GRAY, False, 6),
        ("\u2022  Ejari \u2014 expires Mar 2026  \u26A0 URGENT", RED, True, 11),
        ("\u2022  Emirates ID \u2014 expires May 2026", DARK_GRAY, False, 11),
        ("\u2022  Visa \u2014 expires Aug 2026", DARK_GRAY, False, 11),
        ("Source: Emirates ID, Visa, Ejari documents", MEDIUM_GRAY, False, 9),
    ],
    alignment=PP_ALIGN.LEFT, line_spacing=1.3)

# User follow-up
ub2_y = ai1_y + ai1_h + bubble_gap
ub2 = add_rounded_rect(slide, uc1_x, ub2_y, Inches(3.5), Inches(0.5), LIGHT_BLUE, MEDIUM_BLUE)
add_text_box(slide, uc1_x + Inches(0.15), ub2_y + Inches(0.08), Inches(3.2), Inches(0.35),
             "\U0001F464  \"How do I renew my Ejari?\"",
             font_size=12, color=DARK_BLUE, bold=False)

# AI response 2
ai2_y = ub2_y + Inches(0.6) + bubble_gap
ai2_h = Inches(1.3)
ai2 = add_rounded_rect(slide, uc1_x + Inches(0.5), ai2_y, Inches(5.1), ai2_h, RGBColor(0xF0, 0xF8, 0xF0), GREEN)
add_multiline_text_box(slide, uc1_x + Inches(0.65), ai2_y + Inches(0.08), Inches(4.8), ai2_h - Inches(0.1),
    [
        ("\U0001F916  To renew your Ejari contract:", GREEN, True, 12),
        ("1. Visit Dubai REST app or Ejari website", DARK_GRAY, False, 11),
        ("2. Submit updated tenancy contract", DARK_GRAY, False, 11),
        ("3. Pay renewal fee (AED 220)", DARK_GRAY, False, 11),
        ("Source: Ejari tenancy contract #EJ-2024-XXXXX", MEDIUM_GRAY, False, 9),
    ],
    alignment=PP_ALIGN.LEFT, line_spacing=1.3)


# --- Use Case 2: Instant Document Lookup ---
uc2_x = Inches(6.9)
uc2_w = Inches(5.8)
uc2_y = Inches(1.55)

add_text_box(slide, uc2_x, uc2_y, uc2_w, Inches(0.45),
             "Instant Document Lookup", font_size=18, color=DARK_BLUE, bold=True)

chat_y2 = uc2_y + Inches(0.55)

# User bubble 1
ub3 = add_rounded_rect(slide, uc2_x, chat_y2, Inches(3.8), Inches(0.5), LIGHT_BLUE, MEDIUM_BLUE)
add_text_box(slide, uc2_x + Inches(0.15), chat_y2 + Inches(0.08), Inches(3.5), Inches(0.35),
             "\U0001F464  \"What's my Emirates ID number?\"",
             font_size=12, color=DARK_BLUE, bold=False)

# AI response 1
ai3_y = chat_y2 + Inches(0.6) + bubble_gap
ai3_h = Inches(1.0)
ai3 = add_rounded_rect(slide, uc2_x + Inches(0.5), ai3_y, Inches(5.1), ai3_h, RGBColor(0xF0, 0xF8, 0xF0), GREEN)
add_multiline_text_box(slide, uc2_x + Inches(0.65), ai3_y + Inches(0.08), Inches(4.8), ai3_h - Inches(0.1),
    [
        ("\U0001F916  Your Emirates ID number is:", GREEN, True, 12),
        ("784-XXXX-XXXXXXX-X", DARK_BLUE, True, 14),
        ("Source: Emirates ID, issued 2024-01-15", MEDIUM_GRAY, False, 9),
    ],
    alignment=PP_ALIGN.LEFT, line_spacing=1.4)

# User follow-up
ub4_y = ai3_y + ai3_h + bubble_gap
ub4 = add_rounded_rect(slide, uc2_x, ub4_y, Inches(3.5), Inches(0.5), LIGHT_BLUE, MEDIUM_BLUE)
add_text_box(slide, uc2_x + Inches(0.15), ub4_y + Inches(0.08), Inches(3.2), Inches(0.35),
             "\U0001F464  \"Who is my visa sponsor?\"",
             font_size=12, color=DARK_BLUE, bold=False)

# AI response 2
ai4_y = ub4_y + Inches(0.6) + bubble_gap
ai4_h = Inches(1.3)
ai4 = add_rounded_rect(slide, uc2_x + Inches(0.5), ai4_y, Inches(5.1), ai4_h, RGBColor(0xF0, 0xF8, 0xF0), GREEN)
add_multiline_text_box(slide, uc2_x + Inches(0.65), ai4_y + Inches(0.08), Inches(4.8), ai4_h - Inches(0.1),
    [
        ("\U0001F916  Your visa sponsor details:", GREEN, True, 12),
        ("Sponsor: ACME Technologies LLC", DARK_BLUE, True, 12),
        ("Sponsor ID: 123-XXXX-XXXXXXX-X", DARK_GRAY, False, 11),
        ("Visa Type: Employment Residence", DARK_GRAY, False, 11),
        ("Source: UAE Residence Visa, issued 2024-02-20", MEDIUM_GRAY, False, 9),
    ],
    alignment=PP_ALIGN.LEFT, line_spacing=1.3)

# Vertical separator between the two use cases
add_rect(slide, Inches(6.55), Inches(1.7), Inches(0.02), Inches(5.0), RGBColor(0xDD, 0xDD, 0xDD))

add_slide_number(slide, 5)


# ===================================================================
# SLIDE 6: How It Works (was Slide 5)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "How It Works", "Privacy-preserving RAG architecture")
add_gold_accent_bar(slide)

steps = [
    ("1", "Document\nStorage", "Documents stored\nsecurely in vault\non device", MEDIUM_BLUE),
    ("2", "On-Device\nOCR", "Text extracted\nusing Apple Vision\n/ ML Kit", GREEN),
    ("3", "Embedding\nGeneration", "Text converted to\nvector embeddings\non device", ORANGE),
    ("4", "Semantic\nSearch", "User query matched\nto relevant document\nchunks locally", PURPLE),
    ("5", "Cloud LLM\n(UAE Region)", "Query + top 5 chunks\n(max 2KB) sent to\nAzure OpenAI", DARK_BLUE),
    ("6", "Grounded\nAnswer", "Cited answer returned\nwith source document\nreferences", GOLD),
]

step_w = Inches(1.8)
step_h = Inches(3.2)
start_x = Inches(0.5)
step_y = Inches(1.8)
gap = Inches(0.25)

for i, (num, title, desc, color) in enumerate(steps):
    x = start_x + i * (step_w + gap)

    # Step card
    add_rounded_rect(slide, x, step_y, step_w, step_h, LIGHT_GRAY)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), step_y + Inches(0.15), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_text_box(slide, x + Inches(0.1), step_y + Inches(0.9), step_w - Inches(0.2), Inches(0.7),
                 title, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.2)

    # Description
    add_text_box(slide, x + Inches(0.1), step_y + Inches(1.7), step_w - Inches(0.2), Inches(1.3),
                 desc, font_size=11, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.4)

    # Arrow between steps (cleaner chevron)
    if i < len(steps) - 1:
        arrow_x = x + step_w
        add_text_box(slide, arrow_x - Inches(0.05), step_y + Inches(1.2), Inches(0.35), Inches(0.4),
                     "\u276F", font_size=20, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# On-device vs Cloud divider
add_rect(slide, Inches(0.5), Inches(5.15), Inches(7.8), Inches(0.04), GREEN)
add_text_box(slide, Inches(0.5), Inches(5.02), Inches(3), Inches(0.2),
             "ON-DEVICE (Steps 1\u20134)", font_size=9, color=GREEN, bold=True)

add_rect(slide, Inches(8.7), Inches(5.15), Inches(4.2), Inches(0.04), MEDIUM_BLUE)
add_text_box(slide, Inches(9.5), Inches(5.02), Inches(3), Inches(0.2),
             "CLOUD (Steps 5\u20136)", font_size=9, color=MEDIUM_BLUE, bold=True)

# Privacy banner
banner = add_rect(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.65), DARK_BLUE)
add_text_box(slide, Inches(1.5), Inches(5.55), Inches(10.3), Inches(0.55),
             "\U0001F512  Documents NEVER leave the device  \u2022  Only query + small excerpts sent to UAE-region cloud  \u2022  Zero data retention",
             font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 6)


# ===================================================================
# SLIDE 7: Our Approach (was Slide 6 — Technical Approach)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Our Approach", "Privacy-first architecture for document intelligence")
add_gold_accent_bar(slide)

# --- Option A: Prominent recommended approach ---
opt_a_x = Inches(0.6)
opt_a_y = Inches(1.55)
opt_a_w = Inches(7.5)
opt_a_h = Inches(5.2)

add_rounded_rect(slide, opt_a_x, opt_a_y, opt_a_w, opt_a_h, LIGHT_BLUE, MEDIUM_BLUE)
add_rect(slide, opt_a_x + Inches(0.02), opt_a_y, opt_a_w - Inches(0.04), Inches(0.07), GREEN)

# Recommended badge
badge_w = Inches(2.2)
badge_x = opt_a_x + Inches(0.3)
add_rounded_rect(slide, badge_x, opt_a_y + Inches(0.2), badge_w, Inches(0.4), GREEN)
add_text_box(slide, badge_x, opt_a_y + Inches(0.22), badge_w, Inches(0.35),
             "\u2713  RECOMMENDED", font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, opt_a_x + Inches(0.3), opt_a_y + Inches(0.75), opt_a_w - Inches(0.6), Inches(0.5),
             "Option A: On-Device Processing + Cloud LLM", font_size=22, color=DARK_BLUE, bold=True)

# Key properties in 2-column layout
props_left = [
    ("Privacy", "Documents never leave the device", GREEN),
    ("Data Sent", "Only query + top 5 chunks (max 2KB)", GREEN),
    ("Performance", "2\u20133 second response time", GOLD),
]
props_right = [
    ("Annual Cost", "~$100K / year (1M users)", GREEN),
    ("Arabic", "GPT-4o-mini with strong Arabic support", GREEN),
    ("Offline", "Vector search works offline", GREEN),
]

for j, (label, value, val_color) in enumerate(props_left):
    ry = opt_a_y + Inches(1.5) + j * Inches(0.65)
    add_text_box(slide, opt_a_x + Inches(0.4), ry, Inches(1.2), Inches(0.25),
                 label, font_size=11, color=MEDIUM_GRAY, bold=True)
    add_text_box(slide, opt_a_x + Inches(1.7), ry, Inches(2.5), Inches(0.25),
                 value, font_size=12, color=val_color, bold=False)

for j, (label, value, val_color) in enumerate(props_right):
    ry = opt_a_y + Inches(1.5) + j * Inches(0.65)
    add_text_box(slide, opt_a_x + Inches(4.0), ry, Inches(1.2), Inches(0.25),
                 label, font_size=11, color=MEDIUM_GRAY, bold=True)
    add_text_box(slide, opt_a_x + Inches(5.3), ry, Inches(2.5), Inches(0.25),
                 value, font_size=12, color=val_color, bold=False)

# "Why Option A?" callout box
why_y = opt_a_y + Inches(3.5)
add_rounded_rect(slide, opt_a_x + Inches(0.3), why_y, opt_a_w - Inches(0.6), Inches(1.3), RGBColor(0xE8, 0xF5, 0xE9), GREEN)
add_text_box(slide, opt_a_x + Inches(0.5), why_y + Inches(0.1), opt_a_w - Inches(1.0), Inches(0.35),
             "Why Option A?", font_size=15, color=GREEN, bold=True)
add_text_box(slide, opt_a_x + Inches(0.5), why_y + Inches(0.5), opt_a_w - Inches(1.0), Inches(0.7),
             "\u2022  Privacy-first: documents stay on device, no cloud storage of personal data\n\u2022  Lowest cost: ~$100K/year vs $186K (hybrid) or $450K (full cloud)\n\u2022  Offline capable: core search works without internet",
             font_size=12, color=DARK_GRAY, line_spacing=1.5)

# --- Options B and C: De-emphasized ---
alt_x = Inches(8.5)
alt_w = Inches(4.2)

# Option B
opt_b_y = Inches(1.55)
opt_b_h = Inches(2.4)
add_rounded_rect(slide, alt_x, opt_b_y, alt_w, opt_b_h, LIGHT_GRAY, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, alt_x + Inches(0.2), opt_b_y + Inches(0.15), alt_w - Inches(0.4), Inches(0.35),
             "Option B: Hybrid", font_size=14, color=ORANGE, bold=True)
add_text_box(slide, alt_x + Inches(0.2), opt_b_y + Inches(0.55), alt_w - Inches(0.4), Inches(1.7),
             "\u2022  Embeddings stored in cloud\n\u2022  ~$186K / year\n\u2022  Faster, consistent across devices\n\u2022  Limited offline capability\n\u2022  Higher privacy risk",
             font_size=11, color=MEDIUM_GRAY, line_spacing=1.5)

# Option C
opt_c_y = Inches(4.2)
opt_c_h = Inches(2.55)
add_rounded_rect(slide, alt_x, opt_c_y, alt_w, opt_c_h, LIGHT_GRAY, RGBColor(0xDD, 0xDD, 0xDD))
add_text_box(slide, alt_x + Inches(0.2), opt_c_y + Inches(0.15), alt_w - Inches(0.4), Inches(0.35),
             "Option C: Fully Cloud", font_size=14, color=RED, bold=True)
add_text_box(slide, alt_x + Inches(0.2), opt_c_y + Inches(0.55), alt_w - Inches(0.4), Inches(1.8),
             "\u2022  Full documents uploaded to cloud\n\u2022  ~$450K / year\n\u2022  Most consistent performance\n\u2022  No offline capability\n\u2022  Highest privacy risk\n\u2022  Regulatory concerns",
             font_size=11, color=MEDIUM_GRAY, line_spacing=1.5)

add_slide_number(slide, 7)


# ===================================================================
# SLIDE 8: Phased Rollout (was Slide 7) — No sprint numbers
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Phased Rollout", "Incremental delivery from PoC to General Availability")
add_gold_accent_bar(slide)

phases = [
    ("Phase 0", "Pre-Development", "Internal",
     ["Technical PoC (on-device embedding)", "TDRA policy briefing", "Stakeholder alignment", "User research & consent flows"],
     MEDIUM_BLUE),
    ("Phase 1", "MVP (Official Docs)", "10% users, iOS first",
     ["Emirates ID, Passport, Visa queries", "English + Arabic support", "Citation-backed answers", "In-app feedback collection"],
     GREEN),
    ("Phase 2", "Expansion", "50% users",
     ["Uploaded document support", "Voice input capability", "Android parity", "Performance optimization"],
     ORANGE),
    ("Phase 3", "Advanced + GA", "100% GA",
     ["Document summaries", "Enhanced multilingual", "Export & share answers", "On-device LLM exploration"],
     DARK_BLUE),
]

phase_w = Inches(2.8)
start_x = Inches(0.55)
phase_y = Inches(1.65)
phase_h = Inches(5.1)

for i, (phase, name, rollout, items, color) in enumerate(phases):
    x = start_x + i * (phase_w + Inches(0.25))

    # Card
    add_rounded_rect(slide, x, phase_y, phase_w, phase_h, LIGHT_GRAY)
    # Color header
    add_rect(slide, x + Inches(0.02), phase_y, phase_w - Inches(0.04), Inches(0.95), color)

    # Phase label
    add_text_box(slide, x, phase_y + Inches(0.08), phase_w, Inches(0.35),
                 phase, font_size=13, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, phase_y + Inches(0.38), phase_w, Inches(0.35),
                 name, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Rollout badge
    badge_w = Inches(1.8)
    badge_x = x + (phase_w - badge_w) / 2
    add_rounded_rect(slide, badge_x, phase_y + Inches(1.1), badge_w, Inches(0.35), color)
    add_text_box(slide, badge_x, phase_y + Inches(1.12), badge_w, Inches(0.3),
                 rollout, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Items
    for j, item in enumerate(items):
        iy = phase_y + Inches(1.7) + j * Inches(0.7)
        add_text_box(slide, x + Inches(0.15), iy, phase_w - Inches(0.3), Inches(0.6),
                     f"\u2022  {item}", font_size=12, color=DARK_GRAY, line_spacing=1.3)

    # Arrow connector (except last)
    if i < len(phases) - 1:
        ax = x + phase_w
        add_text_box(slide, ax - Inches(0.05), phase_y + Inches(2.2), Inches(0.35), Inches(0.4),
                     "\u276F", font_size=20, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 8)


# ===================================================================
# SLIDE 9: Privacy & Security (was Slide 8)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Privacy & Security", "Trust-first architecture aligned with UAE data protection standards")
add_gold_accent_bar(slide)

security_items = [
    ("On-Device Processing", "All document OCR, text extraction, and embedding generation\nhappens entirely on the user's device. Documents are never\nuploaded to any external server.", "\U0001F4F1"),
    ("Minimal Cloud Exposure", "Only the user's query and top 5 relevant text chunks\n(maximum 2KB) are sent to the cloud LLM for answer\ngeneration. No full documents are transmitted.", "\u2601"),
    ("Zero Data Retention", "The cloud LLM processes each query statelessly.\nNo user data is stored, cached, or used for model training.\nEach session is ephemeral.", "\U0001F5D1"),
    ("User Consent Required", "Feature is opt-in. Users must explicitly consent before\nfirst use. Clear explanation of what data flows where.\nConsent can be revoked at any time.", "\u2705"),
    ("UAE Data Compliance", "All cloud processing uses Azure OpenAI hosted in\nUAE North region. Compliant with UAE Federal Data\nProtection Law (Federal Decree-Law No. 45 of 2021).", "\U0001F1E6\U0001F1EA"),
]

for i, (title, desc, icon) in enumerate(security_items):
    y = Inches(1.65) + i * Inches(1.1)
    # Row background (alternating)
    if i % 2 == 0:
        add_rect(slide, Inches(0.6), y, Inches(12.1), Inches(1.0), LIGHT_GRAY)

    # Icon area
    icon_bg = add_rounded_rect(slide, Inches(0.8), y + Inches(0.1), Inches(0.7), Inches(0.7), DARK_BLUE)
    add_text_box(slide, Inches(0.8), y + Inches(0.15), Inches(0.7), Inches(0.6),
                 icon, font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide, Inches(1.8), y + Inches(0.05), Inches(3), Inches(0.4),
                 title, font_size=17, color=DARK_BLUE, bold=True)

    # Description
    add_text_box(slide, Inches(1.8), y + Inches(0.4), Inches(10.5), Inches(0.6),
                 desc, font_size=12, color=MEDIUM_GRAY, line_spacing=1.4)

add_slide_number(slide, 9)


# ===================================================================
# SLIDE 10: Success Metrics & Targets (was Slide 9)
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Success Metrics & Targets", "Measurable outcomes across adoption, quality, and performance")
add_gold_accent_bar(slide)

# North star banner
ns_y = Inches(1.6)
add_rounded_rect(slide, Inches(0.6), ns_y, Inches(12.1), Inches(0.8), GOLD)
add_text_box(slide, Inches(0.6), ns_y + Inches(0.05), Inches(12.1), Inches(0.35),
             "NORTH STAR METRIC", font_size=12, color=DARK_BLUE, bold=True,
             alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.6), ns_y + Inches(0.35), Inches(12.1), Inches(0.4),
             "25% of all document access via AI queries by end of Year 1",
             font_size=20, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Metric cards
metrics = [
    ("25%", "User Adoption", "Active users try feature\nwithin 3 months of GA", MEDIUM_BLUE),
    ("5/mo", "Engagement", "Average queries per\nuser per month", GREEN),
    ("95%+", "Accuracy", "Factual accuracy on\nofficial (eSeal) documents", GOLD),
    ("<3s", "Response Time", "90th percentile query\nresponse latency", ORANGE),
    ("4.2+", "Satisfaction", "Star rating in\nin-app feedback", PURPLE),
]

card_w = Inches(2.2)
start_x = Inches(0.55)
card_y = Inches(2.8)
card_h = Inches(2.8)

for i, (value, label, desc, color) in enumerate(metrics):
    x = start_x + i * (card_w + Inches(0.2))

    add_rounded_rect(slide, x, card_y, card_w, card_h, LIGHT_GRAY)
    add_rect(slide, x + Inches(0.02), card_y, card_w - Inches(0.04), Inches(0.06), color)

    # Big number
    add_text_box(slide, x, card_y + Inches(0.3), card_w, Inches(0.7),
                 value, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)

    # Label
    add_text_box(slide, x, card_y + Inches(1.1), card_w, Inches(0.4),
                 label, font_size=16, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

    # Description
    add_text_box(slide, x + Inches(0.15), card_y + Inches(1.6), card_w - Inches(0.3), Inches(1.0),
                 desc, font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER, line_spacing=1.4)

# Phase targets
table_y = Inches(5.9)
add_text_box(slide, Inches(0.8), table_y, Inches(11.7), Inches(0.4),
             "Phase 1: 15% try, 3 q/mo, <5% inaccuracy   |   Phase 2: 25% active, 5 q/mo   |   Phase 3: 30% monthly, 15+ power users",
             font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 10)


# ===================================================================
# SLIDE 11: Next Steps (was Slide 10) — No week/sprint refs
# ===================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_header_bar(slide, "Next Steps", "Immediate actions for Phase 0 kickoff")
add_gold_accent_bar(slide)

next_steps = [
    ("TDRA Buy-in", "Present feature concept and privacy architecture to TDRA for\nstrategic alignment and policy approval.", MEDIUM_BLUE),
    ("DDA Buy-in", "Engage Design Authority for UX direction approval and\ndesign partnership commitment.", PURPLE),
    ("Technical PoC", "Build proof-of-concept: on-device OCR + embedding generation\n+ FAISS vector search + Azure OpenAI UAE North integration.", GREEN),
    ("User Research", "Conduct usability testing with 15\u201320 users on consent flow,\nquery patterns, and answer trustworthiness.", ORANGE),
    ("Go / No-Go Decision", "Evaluate PoC results, user feedback, and stakeholder alignment.\nDecision gate at end of Phase 0.", RED),
]

for i, (title, desc, color) in enumerate(next_steps):
    y = Inches(1.65) + i * Inches(1.08)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.8), y + Inches(0.1),
                                     Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Title
    add_text_box(slide, Inches(1.6), y + Inches(0.0), Inches(5), Inches(0.4),
                 title, font_size=18, color=DARK_BLUE, bold=True)

    # Description
    add_text_box(slide, Inches(1.6), y + Inches(0.4), Inches(10), Inches(0.6),
                 desc, font_size=13, color=MEDIUM_GRAY, line_spacing=1.3)

add_slide_number(slide, 11)


# -- Save --
output_path = r"D:\claude\Ask_Your_Documents_AI_Feature_v2.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"File size: {os.path.getsize(output_path) / 1024:.1f} KB")
print(f"Slides: {len(prs.slides)}")
